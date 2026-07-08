"""
商务汇报 PPT 模板
用法: python business_template.py
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE


class BusinessPPT:
    """商务汇报 PPT 生成器"""

    def __init__(self, title, subtitle="", author="", date=""):
        self.prs = Presentation()
        self.prs.slide_width = Inches(13.333)
        self.prs.slide_height = Inches(7.5)
        self.title = title
        self.subtitle = subtitle
        self.author = author
        self.date = date or "2026年7月"
        # 商务蓝配色
        self.colors = {
            'primary': RGBColor(0x1A, 0x52, 0x76),
            'secondary': RGBColor(0x2E, 0x86, 0xC1),
            'accent': RGBColor(0x85, 0xC1, 0xE9),
            'bg': RGBColor(0xF8, 0xF9, 0xFA),
            'text': RGBColor(0x2C, 0x3E, 0x50),
            'white': RGBColor(0xFF, 0xFF, 0xFF),
            'gray': RGBColor(0x95, 0xA5, 0xA6),
        }

    def add_cover(self):
        """添加封面页"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])

        # 顶部色块
        shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(0), Inches(0),
            Inches(13.333), Inches(3.0)
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = self.colors['primary']
        shape.line.fill.background()

        # 标题
        txBox = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(11), Inches(2))
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        p.text = self.title
        p.font.size = Pt(44)
        p.font.bold = True
        p.font.color.rgb = self.colors['white']
        p.alignment = PP_ALIGN.CENTER

        # 副标题
        if self.subtitle:
            txBox2 = slide.shapes.add_textbox(Inches(1), Inches(4.0), Inches(11), Inches(1.5))
            tf2 = txBox2.text_frame
            p2 = tf2.paragraphs[0]
            p2.text = self.subtitle
            p2.font.size = Pt(22)
            p2.font.color.rgb = self.colors['text']
            p2.alignment = PP_ALIGN.CENTER

        # 底部信息
        txBox3 = slide.shapes.add_textbox(Inches(1), Inches(5.8), Inches(11), Inches(1))
        tf3 = txBox3.text_frame
        p3 = tf3.paragraphs[0]
        p3.text = f"{self.author}  |  {self.date}"
        p3.font.size = Pt(14)
        p3.font.color.rgb = self.colors['gray']
        p3.alignment = PP_ALIGN.CENTER

    def add_section_title(self, section_num, title):
        """添加章节标题页"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])

        # 色块
        shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(0), Inches(2.5),
            Inches(0.6), Inches(2.5)
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = self.colors['primary']
        shape.line.fill.background()

        txBox = slide.shapes.add_textbox(Inches(1.5), Inches(2.5), Inches(10), Inches(1.2))
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        p.text = f"{section_num}"
        p.font.size = Pt(60)
        p.font.bold = True
        p.font.color.rgb = self.colors['primary']

        txBox2 = slide.shapes.add_textbox(Inches(1.5), Inches(3.5), Inches(10), Inches(1))
        tf2 = txBox2.text_frame
        p2 = tf2.paragraphs[0]
        p2.text = title
        p2.font.size = Pt(32)
        p2.font.color.rgb = self.colors['text']

    def add_content_slide(self, title, bullets, note=""):
        """添加内容页"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])

        # 标题栏
        bar = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(0), Inches(0),
            Inches(13.333), Inches(1.2)
        )
        bar.fill.solid()
        bar.fill.fore_color.rgb = self.colors['primary']
        bar.line.fill.background()

        txBox = slide.shapes.add_textbox(Inches(0.8), Inches(0.15), Inches(11), Inches(0.9))
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(28)
        p.font.bold = True
        p.font.color.rgb = self.colors['white']

        # 正文
        txBox2 = slide.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(11), Inches(4.5))
        tf2 = txBox2.text_frame
        tf2.word_wrap = True
        for i, bullet in enumerate(bullets):
            pp = tf2.paragraphs[0] if i == 0 else tf2.add_paragraph()
            pp.text = f"▸ {bullet}"
            pp.font.size = Pt(20)
            pp.font.color.rgb = self.colors['text']
            pp.space_after = Pt(12)

        # 底部备注
        if note:
            txBox3 = slide.shapes.add_textbox(Inches(0.8), Inches(6.3), Inches(11), Inches(0.8))
            tf3 = txBox3.text_frame
            pp = tf3.paragraphs[0]
            pp.text = f"💡 {note}"
            pp.font.size = Pt(12)
            pp.font.italic = True
            pp.font.color.rgb = self.colors['gray']

    def save(self, filename):
        self.prs.save(filename)
        print(f"已生成: {filename}")


if __name__ == "__main__":
    ppt = BusinessPPT(
        title="2026 年度工作汇报",
        subtitle="技术研发中心 · 年度总结与规划",
        author="张三",
        date="2026年7月"
    )

    ppt.add_cover()

    ppt.add_section_title("01", "年度工作回顾")
    ppt.add_content_slide(
        "核心成果",
        ["完成核心系统 V3.0 架构升级，稳定性提升至 99.9%",
         "新产品线 A 完成研发并进入试运行",
         "引入 DevOps 流程，发布效率提升 60%",
         "团队扩充 8 人，关键岗位全部到位"],
        "以上成果均为本年度核心 KPI，已全部达成"
    )

    ppt.add_section_title("02", "重点项目进展")
    ppt.add_content_slide(
        "系统架构升级",
        ["微服务拆分完成 70%，核心模块已上线",
         "数据库读写分离方案已完成 POC",
         "API 网关统一接入，接口响应时间降低 40%"],
        "预计 Q3 完成全部架构升级"
    )

    ppt.add_section_title("03", "下季度规划")
    ppt.add_content_slide(
        "Q3 重点工作",
        ["支付网关升级改造（优先级最高）",
         "数据库读写分离正式上线",
         "微服务架构调研完成并输出方案",
         "持续招聘高级工程师（目标 3-5 人）"]
    )

    ppt.save("商务汇报模板.pptx")
