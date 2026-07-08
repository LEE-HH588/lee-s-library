"""
基础幻灯片生成示例
用法: python basic_ppt.py
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE


def create_basic_ppt():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # ====== 第 1 页：标题页 ======
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # 顶部色块
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0),
        Inches(13.333), Inches(2.5)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0x1A, 0x52, 0x76)
    shape.line.fill.background()

    # 标题
    txBox = slide.shapes.add_textbox(Inches(1), Inches(0.6), Inches(11), Inches(1.5))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "项目汇报"
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    p.alignment = PP_ALIGN.CENTER

    # 副标题
    txBox2 = slide.shapes.add_textbox(Inches(1), Inches(4), Inches(11), Inches(2))
    tf2 = txBox2.text_frame
    p2 = tf2.paragraphs[0]
    p2.text = "2026 年度技术研发中心"
    p2.font.size = Pt(22)
    p2.font.color.rgb = RGBColor(0x66, 0x66, 0x66)
    p2.alignment = PP_ALIGN.CENTER

    p3 = tf2.add_paragraph()
    p3.text = "汇报人：张三  |  2026年7月"
    p3.font.size = Pt(16)
    p3.font.color.rgb = RGBColor(0x99, 0x99, 0x99)
    p3.alignment = PP_ALIGN.CENTER
    p3.space_before = Pt(12)

    # 底部装饰线
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(3), Inches(6.5), Inches(7), Inches(0.04)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = RGBColor(0x1A, 0x52, 0x76)
    line.line.fill.background()

    # ====== 第 2 页：目录页 ======
    slide2 = prs.slides.add_slide(prs.slide_layouts[6])

    txBox = slide2.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(11), Inches(1))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "目 录"
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0x1A, 0x52, 0x76)

    # 目录项
    items = [
        ("01", "年度工作回顾"),
        ("02", "重点项目进展"),
        ("03", "团队建设情况"),
        ("04", "下季度规划"),
    ]
    for i, (num, title) in enumerate(items):
        y = 1.8 + i * 1.2
        # 编号圆
        circle = slide2.shapes.add_shape(
            MSO_SHAPE.OVAL, Inches(1.2), Inches(y), Inches(0.7), Inches(0.7)
        )
        circle.fill.solid()
        circle.fill.fore_color.rgb = RGBColor(0x1A, 0x52, 0x76)
        circle.line.fill.background()
        tf_c = circle.text_frame
        tf_c.paragraphs[0].text = num
        tf_c.paragraphs[0].font.size = Pt(16)
        tf_c.paragraphs[0].font.bold = True
        tf_c.paragraphs[0].font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        tf_c.paragraphs[0].alignment = PP_ALIGN.CENTER

        # 标题文字
        txB = slide2.shapes.add_textbox(Inches(2.3), Inches(y + 0.05), Inches(8), Inches(0.6))
        tf_b = txB.text_frame
        tf_b.paragraphs[0].text = title
        tf_b.paragraphs[0].font.size = Pt(22)
        tf_b.paragraphs[0].font.color.rgb = RGBColor(0x33, 0x33, 0x33)
        tf_b.paragraphs[0].font.bold = True

    # ====== 第 3 页：内容页 ======
    slide3 = prs.slides.add_slide(prs.slide_layouts[6])

    # 标题栏
    bar = slide3.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0),
        Inches(13.333), Inches(1.2)
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = RGBColor(0x1A, 0x52, 0x76)
    bar.line.fill.background()

    txBox = slide3.shapes.add_textbox(Inches(0.8), Inches(0.15), Inches(11), Inches(0.9))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "01  年度工作回顾"
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

    # 三个卡片
    cards = [
        ("核心系统升级", "完成 V3.0 架构升级\n系统稳定性提升至 99.9%"),
        ("新产品研发", "产品线 A 完成研发\n进入试运行阶段"),
        ("流程优化", "引入 DevOps 流程\n发布效率提升 60%"),
    ]
    for i, (title, desc) in enumerate(cards):
        x = 0.8 + i * 4.2
        card = slide3.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(1.8),
            Inches(3.8), Inches(4.5)
        )
        card.fill.solid()
        card.fill.fore_color.rgb = RGBColor(0xF5, 0xF8, 0xFC)
        card.line.color.rgb = RGBColor(0xD5, 0xE3, 0xF0)

        # 卡片标题
        txB = slide3.shapes.add_textbox(Inches(x + 0.3), Inches(2.0), Inches(3.2), Inches(0.8))
        tf_b = txB.text_frame
        tf_b.paragraphs[0].text = title
        tf_b.paragraphs[0].font.size = Pt(20)
        tf_b.paragraphs[0].font.bold = True
        tf_b.paragraphs[0].font.color.rgb = RGBColor(0x1A, 0x52, 0x76)
        tf_b.paragraphs[0].alignment = PP_ALIGN.CENTER

        # 卡片正文
        txB2 = slide3.shapes.add_textbox(Inches(x + 0.3), Inches(3.2), Inches(3.2), Inches(2.5))
        tf_b2 = txB2.text_frame
        tf_b2.word_wrap = True
        for line in desc.split('\n'):
            pp = tf_b2.paragraphs[0] if tf_b2.paragraphs[0].text == '' else tf_b2.add_paragraph()
            pp.text = f"▸ {line}"
            pp.font.size = Pt(16)
            pp.font.color.rgb = RGBColor(0x55, 0x55, 0x55)
            pp.space_after = Pt(8)

    # 保存
    prs.save("基础幻灯片.pptx")
    print("已生成: 基础幻灯片.pptx")


if __name__ == "__main__":
    create_basic_ppt()
