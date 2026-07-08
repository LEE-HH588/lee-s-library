"""
批量生成 PPT 工具 - 从 JSON 数据源自动生成幻灯片
用法: python batch_ppt.py <data.json> [output.pptx]
"""
import json
import sys
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE


def generate_from_json(data_file, output_file="批量生成.pptx"):
    with open(data_file, 'r', encoding='utf-8') as f:
        data = json.load(f)

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # 全局配置
    config = data.get('config', {})
    primary_color = RGBColor(*config.get('primary_color', [0x1A, 0x52, 0x76]))
    title_font_size = config.get('title_font_size', 32)

    for item in data.get('slides', []):
        slide = prs.slides.add_slide(prs.slide_layouts[6])

        # 标题栏
        if item.get('title_bar', True):
            bar = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE, Inches(0), Inches(0),
                Inches(13.333), Inches(1.2)
            )
            bar.fill.solid()
            bar.fill.fore_color.rgb = primary_color
            bar.line.fill.background()

        # 标题
        txBox = slide.shapes.add_textbox(Inches(0.8), Inches(0.15), Inches(11), Inches(0.9))
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        p.text = item.get('title', '')
        p.font.size = Pt(title_font_size)
        p.font.bold = True
        p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

        # 正文（支持多段落 + 子项）
        body_y = item.get('body_y', 1.8)
        txBox2 = slide.shapes.add_textbox(Inches(0.8), Inches(body_y), Inches(11), Inches(5))
        tf2 = txBox2.text_frame
        tf2.word_wrap = True

        first = True
        for line in item.get('body', []):
            if isinstance(line, str):
                if first:
                    pp = tf2.paragraphs[0]
                    first = False
                else:
                    pp = tf2.add_paragraph()
                pp.text = line
                pp.font.size = Pt(item.get('body_font_size', 18))
                pp.font.color.rgb = RGBColor(0x33, 0x33, 0x33)
                pp.space_after = Pt(8)
            elif isinstance(line, list):
                # 子项列表
                for sub in line:
                    pp = tf2.add_paragraph()
                    pp.text = f"• {sub}"
                    pp.font.size = Pt(16)
                    pp.font.color.rgb = RGBColor(0x66, 0x66, 0x66)
                    pp.space_after = Pt(4)
                    pp.level = 1

        # 底部装饰线
        if item.get('bottom_line', True):
            line = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE, Inches(3), Inches(6.8), Inches(7), Inches(0.03)
            )
            line.fill.solid()
            line.fill.fore_color.rgb = primary_color
            line.line.fill.background()

    prs.save(output_file)
    print(f"已生成: {output_file} (共 {len(data.get('slides', []))} 页)")


def generate_sample_data():
    """生成示例 JSON 数据文件"""
    data = {
        "config": {
            "primary_color": [0x1A, 0x52, 0x76],
            "title_font_size": 32
        },
        "slides": [
            {
                "title": "项目进度汇报",
                "body": [
                    "本季度工作整体进展顺利，各项目均按计划推进。",
                    "核心成果：",
                    ["系统架构升级完成，稳定性提升至 99.9%",
                     "新产品线 A 进入试运行阶段",
                     "团队扩充 5 人，新增 2 个重点项目"]
                ]
            },
            {
                "title": "下季度规划",
                "body": [
                    "重点工作方向：",
                    ["推进支付网关升级改造",
                     "启动数据库读写分离方案",
                     "开展微服务架构调研",
                     "持续招募高级开发工程师"],
                    "预计 Q3 末完成全部基础设施升级。"
                ]
            },
            {
                "title": "风险与应对",
                "body": [
                    "当前主要风险：",
                    ["支付网关第三方接口变更 → 已确认兼容方案",
                     "核心人员稳定性 → 已启动人才储备计划",
                     "技术债务累积 → 每迭代预留 20% 重构时间"],
                    "整体风险可控，暂无重大阻塞项。"
                ]
            }
        ]
    }

    with open('sample_data.json', 'w', encoding='utf-8') as f:
        json.dump(data, f, ensure_ascii=False, indent=2)
    print("已生成示例数据: sample_data.json")


if __name__ == "__main__":
    if len(sys.argv) < 2:
        print("用法: python batch_ppt.py <data.json> [output.pptx]")
        print("示例: python batch_ppt.py sample_data.pptx")
        print()
        print("首先生成示例数据:")
        generate_sample_data()
    else:
        data_file = sys.argv[1]
        output_file = sys.argv[2] if len(sys.argv) > 2 else "批量生成.pptx"
        generate_from_json(data_file, output_file)
