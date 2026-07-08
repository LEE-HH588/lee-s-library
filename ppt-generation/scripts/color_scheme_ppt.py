"""
专业配色方案幻灯片示例
用法: python color_scheme_ppt.py
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE


# ======== 预设配色方案 ========
COLOR_SCHEMES = {
    'business_blue': {
        'name': '商务蓝',
        'primary': RGBColor(0x1A, 0x52, 0x76),
        'secondary': RGBColor(0x2E, 0x86, 0xC1),
        'accent': RGBColor(0x85, 0xC1, 0xE9),
        'bg': RGBColor(0xF8, 0xF9, 0xFA),
        'text': RGBColor(0x2C, 0x3E, 0x50),
        'muted': RGBColor(0x7F, 0x8C, 0x8D),
        'gradient_start': RGBColor(0x1A, 0x52, 0x76),
        'gradient_end': RGBColor(0x85, 0xC1, 0xE9),
    },
    'tech_blue': {
        'name': '科技蓝',
        'primary': RGBColor(0x0F, 0x17, 0x2A),
        'secondary': RGBColor(0x3B, 0x82, 0xF6),
        'accent': RGBColor(0x60, 0xA5, 0xFA),
        'bg': RGBColor(0xF0, 0xF9, 0xFF),
        'text': RGBColor(0x1E, 0x29, 0x3B),
        'muted': RGBColor(0x64, 0x74, 0x8B),
        'gradient_start': RGBColor(0x0F, 0x17, 0x2A),
        'gradient_end': RGBColor(0x3B, 0x82, 0xF6),
    },
    'nature_green': {
        'name': '自然绿',
        'primary': RGBColor(0x1E, 0x84, 0x49),
        'secondary': RGBColor(0x27, 0xAE, 0x60),
        'accent': RGBColor(0x82, 0xE0, 0xAA),
        'bg': RGBColor(0xF0, 0xFF, 0xF0),
        'text': RGBColor(0x1C, 0x3D, 0x2A),
        'muted': RGBColor(0x6B, 0x8E, 0x76),
        'gradient_start': RGBColor(0x1E, 0x84, 0x49),
        'gradient_end': RGBColor(0x82, 0xE0, 0xAA),
    },
    'warm_orange': {
        'name': '热情橙',
        'primary': RGBColor(0xD3, 0x54, 0x00),
        'secondary': RGBColor(0xE6, 0x7E, 0x22),
        'accent': RGBColor(0xFA, 0xD7, 0xA1),
        'bg': RGBColor(0xFF, 0xFB, 0xF5),
        'text': RGBColor(0x5D, 0x3A, 0x1A),
        'muted': RGBColor(0xA0, 0x7A, 0x5A),
        'gradient_start': RGBColor(0xD3, 0x54, 0x00),
        'gradient_end': RGBColor(0xFA, 0xD7, 0xA1),
    },
    'elegant_purple': {
        'name': '典雅紫',
        'primary': RGBColor(0x6C, 0x34, 0x83),
        'secondary': RGBColor(0x8E, 0x44, 0xAD),
        'accent': RGBColor(0xD2, 0xB4, 0xDE),
        'bg': RGBColor(0xFD, 0xF8, 0xFF),
        'text': RGBColor(0x3A, 0x1D, 0x4A),
        'muted': RGBColor(0x8E, 0x7A, 0x9A),
        'gradient_start': RGBColor(0x6C, 0x34, 0x83),
        'gradient_end': RGBColor(0xD2, 0xB4, 0xDE),
    },
}


def apply_gradient(slide, color1, color2, angle=90):
    """为幻灯片设置渐变色背景"""
    bg = slide.background
    fill = bg.fill
    fill.gradient()
    fill.gradient.angle = angle
    stops = fill.gradient.stops
    stops[0].position = 0.0
    stops[0].color.rgb = color1
    stops[1].position = 1.0
    stops[1].color.rgb = color2


def create_color_slide(prs, scheme_key):
    """创建单个配色方案展示页"""
    c = COLOR_SCHEMES[scheme_key]
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # 渐变色背景
    apply_gradient(slide, c['gradient_start'], c['bg'], angle=90)

    # 标题
    txBox = slide.shapes.add_textbox(Inches(0.8), Inches(0.3), Inches(11), Inches(1))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = f"配色方案：{c['name']}"
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = c['primary']

    # 色卡展示
    colors_info = [
        ('primary', '主色', c['primary']),
        ('secondary', '辅色', c['secondary']),
        ('accent', '强调色', c['accent']),
        ('bg', '背景色', c['bg']),
        ('text', '文字色', c['text']),
        ('muted', '弱化色', c['muted']),
    ]

    for i, (key, label, color) in enumerate(colors_info):
        x = 1.0 + (i % 3) * 4.0
        y = 1.8 + (i // 3) * 3.0

        # 色块
        swatch = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(x), Inches(y), Inches(3.2), Inches(2.0)
        )
        swatch.fill.solid()
        swatch.fill.fore_color.rgb = color
        swatch.line.color.rgb = RGBColor(0xDD, 0xDD, 0xDD)

        # 标签
        txB = slide.shapes.add_textbox(Inches(x + 0.2), Inches(y + 0.3), Inches(2.8), Inches(0.5))
        tf_b = txB.text_frame
        tf_b.paragraphs[0].text = label
        tf_b.paragraphs[0].font.size = Pt(14)
        tf_b.paragraphs[0].font.bold = True
        tf_b.paragraphs[0].font.color.rgb = (
            RGBColor(0xFF, 0xFF, 0xFF) if key in ('primary', 'secondary', 'text')
            else RGBColor(0x33, 0x33, 0x33)
        )
        tf_b.paragraphs[0].alignment = PP_ALIGN.CENTER

        # 颜色值
        txB2 = slide.shapes.add_textbox(Inches(x + 0.2), Inches(y + 0.8), Inches(2.8), Inches(0.5))
        tf_b2 = txB2.text_frame
        hex_val = '#{:02X}{:02X}{:02X}'.format(color[0], color[1], color[2])
        tf_b2.paragraphs[0].text = hex_val
        tf_b2.paragraphs[0].font.size = Pt(12)
        tf_b2.paragraphs[0].font.color.rgb = (
            RGBColor(0xFF, 0xFF, 0xFF) if key in ('primary', 'secondary', 'text')
            else RGBColor(0x66, 0x66, 0x66)
        )
        tf_b2.paragraphs[0].alignment = PP_ALIGN.CENTER


def create_color_ppt():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # 为每个配色方案创建一页
    for key in COLOR_SCHEMES:
        create_color_slide(prs, key)

    prs.save("配色方案展示.pptx")
    print("已生成: 配色方案展示.pptx (共 5 页)")


if __name__ == "__main__":
    create_color_ppt()
