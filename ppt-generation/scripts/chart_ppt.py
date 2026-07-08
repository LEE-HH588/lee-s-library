"""
数据图表幻灯片生成示例
用法: python chart_ppt.py
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION, XL_LABEL_POSITION
from pptx.chart.data import CategoryChartData
from pptx.enum.shapes import MSO_SHAPE


def create_chart_ppt():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # ====== 封面 ======
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0),
        Inches(13.333), Inches(7.5)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0x0F, 0x17, 0x2A)
    shape.line.fill.background()

    txBox = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(11), Inches(3))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "数据报告"
    p.font.size = Pt(48)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    p.alignment = PP_ALIGN.CENTER

    p2 = tf.add_paragraph()
    p2.text = "2026 年第二季度 · 销售数据分析"
    p2.font.size = Pt(20)
    p2.font.color.rgb = RGBColor(0x60, 0xA5, 0xFA)
    p2.alignment = PP_ALIGN.CENTER
    p2.space_before = Pt(16)

    # ====== 柱状图页 ======
    slide2 = prs.slides.add_slide(prs.slide_layouts[6])

    txBox = slide2.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.8))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "各区域销售额统计"
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0x0F, 0x17, 0x2A)

    p2 = tf.add_paragraph()
    p2.text = "单位：万元"
    p2.font.size = Pt(12)
    p2.font.color.rgb = RGBColor(0x99, 0x99, 0x99)

    # 柱状图数据
    chart_data = CategoryChartData()
    chart_data.categories = ['华东区', '华南区', '华北区', '西南区', '西北区']
    chart_data.add_series('Q1', (1280, 950, 720, 430, 280))
    chart_data.add_series('Q2', (1560, 1120, 890, 550, 360))

    chart_frame = slide2.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED,
        Inches(0.5), Inches(1.5), Inches(8), Inches(5.5),
        chart_data
    )

    chart = chart_frame.chart
    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.BOTTOM
    chart.legend.include_in_layout = False

    # 系列颜色
    plot = chart.plots[0]
    plot.gap_width = 80
    series0 = plot.series[0]
    series0.format.fill.solid()
    series0.format.fill.fore_color.rgb = RGBColor(0x60, 0xA5, 0xFA)

    series1 = plot.series[1]
    series1.format.fill.solid()
    series1.format.fill.fore_color.rgb = RGBColor(0x3B, 0x82, 0xF6)

    # 数据标签
    for series in [series0, series1]:
        series.has_data_labels = True
        series.data_labels.font.size = Pt(10)
        series.data_labels.font.color.rgb = RGBColor(0x33, 0x33, 0x33)
        series.data_labels.number_format = '#,##0'

    # 右侧说明卡片
    card = slide2.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(9), Inches(1.5), Inches(4), Inches(3.5)
    )
    card.fill.solid()
    card.fill.fore_color.rgb = RGBColor(0xF0, 0xF9, 0xFF)
    card.line.color.rgb = RGBColor(0x60, 0xA5, 0xFA)

    txB = slide2.shapes.add_textbox(Inches(9.3), Inches(1.7), Inches(3.5), Inches(3))
    tf_b = txB.text_frame
    tf_b.word_wrap = True
    insights = [
        "📊 数据洞察",
        "",
        "▸ 华东区 Q2 销售额 1,560 万，排名第一",
        "▸ 华南区环比增长 +17.9%",
        "▸ 西南区增长最快 (+27.9%)",
        "▸ 整体环比增长 +21.9%",
    ]
    for i, line in enumerate(insights):
        if i == 0:
            p = tf_b.paragraphs[0]
            p.text = line
            p.font.size = Pt(14)
            p.font.bold = True
            p.font.color.rgb = RGBColor(0x0F, 0x17, 0x2A)
        else:
            pp = tf_b.add_paragraph()
            pp.text = line
            pp.font.size = Pt(12)
            pp.font.color.rgb = RGBColor(0x55, 0x55, 0x55)
            pp.space_after = Pt(4)

    # ====== 饼图页 ======
    slide3 = prs.slides.add_slide(prs.slide_layouts[6])

    txBox = slide3.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.8))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "市场份额分布"
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0x0F, 0x17, 0x2A)

    pie_data = CategoryChartData()
    pie_data.categories = ['华东区', '华南区', '华北区', '西南区', '西北区']
    pie_data.add_series('市场份额', (37.8, 27.2, 21.6, 13.3, 8.7))

    chart_frame2 = slide3.shapes.add_chart(
        XL_CHART_TYPE.PIE,
        Inches(0.5), Inches(1.5), Inches(7), Inches(5.5),
        pie_data
    )

    chart2 = chart_frame2.chart
    chart2.has_legend = True
    chart2.legend.position = XL_LEGEND_POSITION.BOTTOM
    chart2.legend.font.size = Pt(12)

    # 饼图颜色
    pie_colors = [
        RGBColor(0x3B, 0x82, 0xF6),  # 蓝
        RGBColor(0x10, 0xB9, 0x81),  # 绿
        RGBColor(0xF5, 0x9E, 0x0B),  # 橙
        RGBColor(0xEF, 0x44, 0x44),  # 红
        RGBColor(0x8B, 0x5C, 0xF6),  # 紫
    ]
    plot2 = chart2.plots[0]
    for i, color in enumerate(pie_colors):
        point = plot2.series[0].points[i]
        point.format.fill.solid()
        point.format.fill.fore_color.rgb = color

    # 数据标签
    series_pie = plot2.series[0]
    series_pie.has_data_labels = True
    series_pie.data_labels.show_percentage = True
    series_pie.data_labels.show_category_name = True
    series_pie.data_labels.font.size = Pt(10)

    # 保存
    prs.save("数据图表幻灯片.pptx")
    print("已生成: 数据图表幻灯片.pptx")


if __name__ == "__main__":
    create_chart_ppt()
