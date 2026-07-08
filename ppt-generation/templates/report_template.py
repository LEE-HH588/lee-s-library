"""
数据报告 PPT 模板
用法: python report_template.py
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.chart.data import CategoryChartData
from pptx.enum.shapes import MSO_SHAPE


class ReportPPT:
    """数据报告 PPT 生成器"""

    def __init__(self, title, subtitle=""):
        self.prs = Presentation()
        self.prs.slide_width = Inches(13.333)
        self.prs.slide_height = Inches(7.5)
        self.title = title
        self.subtitle = subtitle
        # 科技蓝配色
        self.c = {
            'primary': RGBColor(0x0F, 0x17, 0x2A),
            'secondary': RGBColor(0x3B, 0x82, 0xF6),
            'accent': RGBColor(0x60, 0xA5, 0xFA),
            'bg': RGBColor(0xF0, 0xF9, 0xFF),
            'text': RGBColor(0x1E, 0x29, 0x3B),
            'white': RGBColor(0xFF, 0xFF, 0xFF),
            'light': RGBColor(0xDD, 0xDD, 0xDD),
            'success': RGBColor(0x10, 0xB9, 0x81),
            'warning': RGBColor(0xF5, 0x9E, 0x0B),
        }

    def add_cover(self):
        """添加封面目（深色背景）"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])

        bg = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(0), Inches(0),
            Inches(13.333), Inches(7.5)
        )
        bg.fill.solid()
        bg.fill.fore_color.rgb = self.c['primary']
        bg.line.fill.background()

        txBox = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(11), Inches(3))
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        p.text = self.title
        p.font.size = Pt(44)
        p.font.bold = True
        p.font.color.rgb = self.c['white']
        p.alignment = PP_ALIGN.CENTER

        if self.subtitle:
            p2 = tf.add_paragraph()
            p2.text = self.subtitle
            p2.font.size = Pt(18)
            p2.font.color.rgb = self.c['accent']
            p2.alignment = PP_ALIGN.CENTER
            p2.space_before = Pt(16)

    def add_chart_slide(self, title, categories, series_data, chart_type="bar"):
        """添加图表页"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])

        # 左侧标题
        txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.8))
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(28)
        p.font.bold = True
        p.font.color.rgb = self.c['primary']

        # 图表数据
        chart_data = CategoryChartData()
        chart_data.categories = categories
        for name, values in series_data:
            chart_data.add_series(name, values)

        # 图表类型
        chart_type_map = {
            'bar': XL_CHART_TYPE.COLUMN_CLUSTERED,
            'line': XL_CHART_TYPE.LINE_MARKERS,
            'pie': XL_CHART_TYPE.PIE,
        }
        xl_type = chart_type_map.get(chart_type, XL_CHART_TYPE.COLUMN_CLUSTERED)

        chart_frame = slide.shapes.add_chart(
            xl_type, Inches(0.5), Inches(1.5),
            Inches(8), Inches(5.5), chart_data
        )

        chart = chart_frame.chart
        chart.has_legend = True
        chart.legend.position = XL_LEGEND_POSITION.BOTTOM

        # 设置颜色
        colors = [self.c['secondary'], self.c['success'],
                  self.c['warning'], self.c['accent']]
        plot = chart.plots[0]
        for i in range(len(series_data)):
            s = plot.series[i]
            s.format.fill.solid()
            s.format.fill.fore_color.rgb = colors[i % len(colors)]

        return slide

    def add_kpi_slide(self, title, kpi_cards):
        """添加 KPI 指标页"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])

        txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.8))
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(28)
        p.font.bold = True
        p.font.color.rgb = self.c['primary']

        for i, (label, value, change) in enumerate(kpi_cards):
            x = 0.8 + i * 3.2
            card = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(x), Inches(1.8), Inches(2.8), Inches(3.5)
            )
            card.fill.solid()
            card.fill.fore_color.rgb = self.c['bg']
            card.line.color.rgb = self.c['light']

            # 数值
            txB = slide.shapes.add_textbox(Inches(x), Inches(2.2), Inches(2.8), Inches(1))
            tf_b = txB.text_frame
            tf_b.paragraphs[0].text = value
            tf_b.paragraphs[0].font.size = Pt(36)
            tf_b.paragraphs[0].font.bold = True
            tf_b.paragraphs[0].font.color.rgb = self.c['secondary']
            tf_b.paragraphs[0].alignment = PP_ALIGN.CENTER

            # 变化
            txB2 = slide.shapes.add_textbox(Inches(x), Inches(3.2), Inches(2.8), Inches(0.5))
            tf_b2 = txB2.text_frame
            tf_b2.paragraphs[0].text = change
            tf_b2.paragraphs[0].font.size = Pt(14)
            tf_b2.paragraphs[0].font.color.rgb = (
                self.c['success'] if change.startswith('+') else self.c['text']
            )
            tf_b2.paragraphs[0].alignment = PP_ALIGN.CENTER

            # 标签
            txB3 = slide.shapes.add_textbox(Inches(x), Inches(3.8), Inches(2.8), Inches(0.5))
            tf_b3 = txB3.text_frame
            tf_b3.paragraphs[0].text = label
            tf_b3.paragraphs[0].font.size = Pt(14)
            tf_b3.paragraphs[0].font.color.rgb = self.c['text']
            tf_b3.paragraphs[0].alignment = PP_ALIGN.CENTER

    def save(self, filename):
        self.prs.save(filename)
        print(f"已生成: {filename}")


if __name__ == "__main__":
    ppt = ReportPPT(
        title="Q2 数据报告",
        subtitle="2026 年第二季度 · 销售运营数据分析"
    )

    ppt.add_cover()

    ppt.add_kpi_slide("核心指标概览", [
        ("总销售额", "¥4,120 万", "+21.9%"),
        ("订单量", "12,580 单", "+15.3%"),
        ("客单价", "¥3,275", "+5.7%"),
        ("客户数", "8,920", "+18.1%"),
    ])

    ppt.add_chart_slide(
        "各区域销售趋势",
        ['华东', '华南', '华北', '西南', '西北'],
        [('Q1', (1280, 950, 720, 430, 280)),
         ('Q2', (1560, 1120, 890, 550, 360))],
        chart_type='bar'
    )

    ppt.save("数据报告模板.pptx")
